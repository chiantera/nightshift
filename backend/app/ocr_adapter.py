from __future__ import annotations

import base64
import io
import os
from abc import ABC, abstractmethod

from .ocr_models import OcrInput, OcrPage, OcrResult, OcrWarning


class OcrAdapter(ABC):
    """Stable boundary between case analysis and replaceable OCR engines."""

    engine: str

    @abstractmethod
    def extract(self, request: OcrInput) -> OcrResult:
        """Return normalized OCR text, page refs, confidence, and warnings."""


class TextOnlyOcrAdapter(OcrAdapter):
    """Alpha contract stub.

    This intentionally does not run OCR. It only passes through text/plain input
    so downstream code can integrate against the normalized OCR contract before
    DeepSeek-OCR-2, Rust, or hosted engines are selected.
    """

    engine = "text-only-placeholder"

    def extract(self, request: OcrInput) -> OcrResult:
        if request.mime_type != "text/plain":
            return OcrResult(
                success=False,
                engine=self.engine,
                pages=[],
                warnings=[
                    OcrWarning(
                        code="unsupported_mime_type",
                        message=(
                            "text-only-placeholder only supports text/plain "
                            "in the alpha contract stub."
                        ),
                    )
                ],
            )

        raw_content = request.content
        if raw_content is None and request.file_path is not None:
            raw_content = request.file_path.read_bytes()

        text = (raw_content or b"").decode("utf-8")
        return OcrResult(
            success=True,
            engine=self.engine,
            pages=[OcrPage(page=1, text=text, confidence=1.0, blocks=[])],
            warnings=[
                OcrWarning(
                    code="placeholder_engine",
                    message="Text/plain passthrough only; no OCR inference was run.",
                    severity="warning",
                )
            ],
        )


class PypdfAdapter(OcrAdapter):
    """Extract text from native (non-scanned) PDFs using pypdf — free, instant, local.

    Returns success=False for non-PDF input or PDFs with no extractable text
    (e.g. scanned-only), so the caller can fall back to MistralOcrAdapter.
    """

    engine = "pypdf"

    def extract(self, request: OcrInput) -> OcrResult:
        if request.mime_type != "application/pdf":
            return OcrResult(
                success=False, engine=self.engine, pages=[],
                warnings=[OcrWarning(code="unsupported_mime_type", message=f"pypdf only handles application/pdf, got {request.mime_type}")],
            )

        try:
            import pypdf  # noqa: PLC0415
        except ImportError:
            return OcrResult(
                success=False, engine=self.engine, pages=[],
                warnings=[OcrWarning(code="missing_dependency", message="pypdf not installed", severity="error")],
            )

        content = request.content
        if content is None and request.file_path is not None:
            content = request.file_path.read_bytes()
        if not content:
            return OcrResult(success=False, engine=self.engine, pages=[],
                             warnings=[OcrWarning(code="empty_content", message="No content provided")])

        reader = pypdf.PdfReader(io.BytesIO(content))
        pages: list[OcrPage] = []
        for i, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                pages.append(OcrPage(page=i, text=text, confidence=1.0, blocks=[]))

        if not pages:
            return OcrResult(
                success=False, engine=self.engine, pages=[],
                warnings=[OcrWarning(code="no_text_layer", message="PDF has no extractable text layer — likely scanned. Try MistralOcrAdapter.")],
            )

        return OcrResult(success=True, engine=self.engine, pages=pages, warnings=[])


class PptxAdapter(OcrAdapter):
    """Extract text from PPTX files using python-pptx."""

    engine = "python-pptx"

    def extract(self, request: OcrInput) -> OcrResult:
        SUPPORTED = [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ]
        if request.mime_type not in SUPPORTED and not any(request.mime_type.endswith(s) for s in ["presentationml.presentation", ".powerpoint"]):
            return OcrResult(
                success=False, engine=self.engine, pages=[],
                warnings=[OcrWarning(code="unsupported_mime_type", message=f"python-pptx only handles PPT/PPTX, got {request.mime_type}")],
            )

        try:
            from pptx import Presentation  # noqa: PLC0415
        except ImportError:
            return OcrResult(
                success=False, engine=self.engine, pages=[],
                warnings=[OcrWarning(code="missing_dependency", message="python-pptx not installed", severity="error")],
            )

        content = request.content
        if content is None and request.file_path is not None:
            content = request.file_path.read_bytes()
        if not content:
            return OcrResult(success=False, engine=self.engine, pages=[],
                             warnings=[OcrWarning(code="empty_content", message="No content provided")])

        prs = Presentation(io.BytesIO(content))
        pages: list[OcrPage] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        t = paragraph.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            t = cell.text.strip()
                            if t:
                                texts.append(t)
                if shape.has_chart:
                    texts.append(f"[Grafico: {shape.chart.chart_type}]")
            if texts:
                pages.append(OcrPage(page=i, text="\n".join(texts), confidence=0.9, blocks=[]))

        if not pages:
            return OcrResult(
                success=False, engine=self.engine, pages=[],
                warnings=[OcrWarning(code="no_text_layer", message="PPTX has no extractable text — likely image-only slides. Convert to PDF for OCR.")],
            )

        return OcrResult(success=True, engine=self.engine, pages=pages, warnings=[])


class XlsxAdapter(OcrAdapter):
    """Extract text from XLSX files using openpyxl."""

    engine = "python-openpyxl"

    def extract(self, request: OcrInput) -> OcrResult:
        SUPPORTED = [
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ]
        if request.mime_type not in SUPPORTED and not any(request.mime_type.endswith(s) for s in ["spreadsheetml.sheet", ".excel"]):
            return OcrResult(
                success=False, engine=self.engine, pages=[],
                warnings=[OcrWarning(code="unsupported_mime_type", message=f"openpyxl only handles XLS/XLSX, got {request.mime_type}")],
            )

        try:
            import openpyxl  # noqa: PLC0415
        except ImportError:
            return OcrResult(
                success=False, engine=self.engine, pages=[],
                warnings=[OcrWarning(code="missing_dependency", message="openpyxl not installed", severity="error")],
            )

        content = request.content
        if content is None and request.file_path is not None:
            content = request.file_path.read_bytes()
        if not content:
            return OcrResult(success=False, engine=self.engine, pages=[],
                             warnings=[OcrWarning(code="empty_content", message="No content provided")])

        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        pages: list[OcrPage] = []
        for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=1):
            ws = wb[sheet_name]
            rows_text: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows_text.append(" | ".join(cells))
            if rows_text:
                page_title = f"Foglio: {sheet_name}"
                pages.append(OcrPage(page=sheet_idx, text=page_title + "\n" + "\n".join(rows_text), confidence=0.9, blocks=[]))

        wb.close()

        if not pages:
            return OcrResult(
                success=False, engine=self.engine, pages=[],
                warnings=[OcrWarning(code="empty_result", message="XLSX has no readable data cells")],
            )

        return OcrResult(success=True, engine=self.engine, pages=pages, warnings=[])


class MistralOcrAdapter(OcrAdapter):
    """OCR via Mistral OCR API (mistral-ocr-latest).

    Handles scanned PDFs and images. Requires MISTRAL_API_KEY env var.
    Returns markdown-formatted text per page.
    """

    engine = "mistral-ocr"

    def extract(self, request: OcrInput) -> OcrResult:
        api_key = os.environ.get("MISTRAL_API_KEY", "")
        if not api_key:
            return OcrResult(
                success=False, engine=self.engine, pages=[],
                warnings=[OcrWarning(code="missing_api_key", message="MISTRAL_API_KEY not set", severity="error")],
            )

        try:
            from mistralai.client import Mistral  # noqa: PLC0415
        except ImportError:
            return OcrResult(
                success=False, engine=self.engine, pages=[],
                warnings=[OcrWarning(code="missing_dependency", message="mistralai not installed", severity="error")],
            )

        content = request.content
        if content is None and request.file_path is not None:
            content = request.file_path.read_bytes()
        if not content:
            return OcrResult(success=False, engine=self.engine, pages=[],
                             warnings=[OcrWarning(code="empty_content", message="No content provided")])

        is_pdf = request.mime_type == "application/pdf"
        doc_type = "document_url" if is_pdf else "image_url"
        b64 = base64.standard_b64encode(content).decode()
        data_url = f"data:{request.mime_type};base64,{b64}"

        client = Mistral(api_key=api_key)
        response = client.ocr.process(
            model="mistral-ocr-latest",
            document={"type": doc_type, doc_type: data_url},
        )

        pages: list[OcrPage] = []
        for p in response.pages:
            text = (p.markdown or "").strip()
            if text:
                pages.append(OcrPage(page=p.index + 1, text=text, confidence=0.95, blocks=[]))

        if not pages:
            return OcrResult(success=False, engine=self.engine, pages=[],
                             warnings=[OcrWarning(code="empty_result", message="Mistral OCR returned no text")])

        return OcrResult(success=True, engine=self.engine, pages=pages, warnings=[])
